import streamlit as st
import json
from pptx import Presentation
from pptx.util import Inches, Pt
from pptx.dml.color import RGBColor
from pptx.enum.shapes import MSO_SHAPE
import tempfile
from dotenv import load_dotenv
import os
from openai import OpenAI

# ─── Security Check (Passkey via URL) ────────────────────────────────────────
query_params = st.query_params
user_key = query_params.get("key", None)
required_key = st.secrets["security"]["access_key"]

if user_key != required_key:
    st.error("🔒 Access Denied: Invalid or missing key in URL.")
    st.stop()

# Load .env
load_dotenv()
client = OpenAI(api_key=st.secrets["OPENAI_API_KEY"])

st.set_page_config(page_title="Modern AI PPT Enhancer", layout="centered")
st.title("🎨 AI-Enhanced PowerPoint Generator (Modern Style)")

# Extract content from uploaded PPTX
def extract_ppt_text(ppt_file):
    prs = Presentation(ppt_file)
    content = []
    for slide in prs.slides:
        slide_data = {"heading": "", "bullet_points": []}
        if slide.shapes.title:
            slide_data["heading"] = slide.shapes.title.text.strip()
        for shape in slide.shapes:
            if shape.has_text_frame and shape != slide.shapes.title:
                for para in shape.text_frame.paragraphs:
                    text = para.text.strip()
                    if text:
                        slide_data["bullet_points"].append(text)
        content.append(slide_data)
    return content

# OpenAI GPT-4 prompt to enhance structure
def generate_improved_json(topic, extracted, slide_count, style):
    system = "You are a presentation designer. Improve structure of presentation slides and return clean JSON."
    user = f"""
Topic: {topic}
Slides to improve:
{json.dumps(extracted, indent=2)}

Instructions:
- Target slide count: ~{slide_count}
- Style: {style}
- Add chart_type (bar, pie, line) when appropriate
- Return only valid JSON like:
{{
  "title": "Title",
  "slides": [
    {{
      "heading": "Slide title",
      "bullet_points": ["...", "..."],
      "chart_type": "bar"  # optional
    }}
  ]
}}
"""

    response = client.chat.completions.create(
        model="gpt-4",
        messages=[
            {"role": "system", "content": system},
            {"role": "user", "content": user}
        ],
        temperature=0.7,
        max_tokens=1200
    )

    return response.choices[0].message.content.strip()

# Generate styled presentation
def json_to_ppt(content, output_path):
    prs = Presentation()
    blank_layout = prs.slide_layouts[6]  # fully blank

    # Title slide
    title_slide = prs.slides.add_slide(prs.slide_layouts[0])
    title_slide.shapes.title.text = content.get("title", "Untitled Presentation")
    title_slide.placeholders[1].text = "Enhanced by GPT-4"

    for s in content.get("slides", []):
        slide = prs.slides.add_slide(blank_layout)

        # Background color (light blue)
        background = slide.background
        fill = background.fill
        fill.solid()
        fill.fore_color.rgb = RGBColor(240, 248, 255)

        # Title shape (dark blue box with white text)
        title_shape = slide.shapes.add_shape(
            MSO_SHAPE.ROUNDED_RECTANGLE, Inches(0.5), Inches(0.4), Inches(9), Inches(1)
        )
        title_shape.fill.solid()
        title_shape.fill.fore_color.rgb = RGBColor(0, 102, 204)
        title_shape.text_frame.text = s.get("heading", "")
        title_p = title_shape.text_frame.paragraphs[0]
        title_p.font.size = Pt(28)
        title_p.font.color.rgb = RGBColor(255, 255, 255)
        title_p.font.bold = True

        # Bullet points box
        textbox = slide.shapes.add_textbox(Inches(0.8), Inches(1.6), Inches(8.5), Inches(4.5))
        tf = textbox.text_frame
        tf.word_wrap = True

        for bp in s.get("bullet_points", []):
            p = tf.add_paragraph()
            p.text = bp
            p.font.size = Pt(18)
            p.level = 0  # 0 to 4 for indentation

        # Chart placeholder shape (optional)
        if s.get("chart_type"):
            chart_box = slide.shapes.add_shape(
                MSO_SHAPE.RECTANGLE, Inches(5.5), Inches(5), Inches(3.5), Inches(1)
            )
            chart_box.fill.solid()
            chart_box.fill.fore_color.rgb = RGBColor(230, 230, 230)
            chart_box.text = f"[{s['chart_type']} chart placeholder]"
            chart_box.text_frame.paragraphs[0].font.size = Pt(14)
            chart_box.text_frame.paragraphs[0].font.color.rgb = RGBColor(80, 80, 80)

    prs.save(output_path)
    return output_path

# Streamlit form
with st.form("ppt_form"):
    st.subheader("Upload & Enhance Presentation")
    ppt_file = st.file_uploader("📤 Upload a .pptx file", type=["pptx"])
    topic = st.text_input("🧠 Topic", "AI in Healthcare")
    slide_count = st.slider("🧾 Desired Slide Count", 3, 20, 6)
    style = st.text_input("🎨 Style", "Professional, modern, minimal")

    submitted = st.form_submit_button("✨ Enhance Presentation")

if submitted:
    if not ppt_file:
        st.error("Please upload a presentation file.")
    else:
        with st.spinner("🔍 Extracting slides..."):
            extracted = extract_ppt_text(ppt_file)
            st.success("✅ Slide content extracted")
            st.json(extracted)

        with st.spinner("🤖 Generating improved version via GPT-4..."):
            try:
                json_str = generate_improved_json(topic, extracted, slide_count, style)
                data = json.loads(json_str)
                st.success("✅ JSON structure ready")
                st.json(data)

                with st.spinner("🛠️ Building modern styled PPTX..."):
                    with tempfile.NamedTemporaryFile(delete=False, suffix=".pptx") as tmp:
                        pptx_path = json_to_ppt(data, tmp.name)
                        with open(pptx_path, "rb") as f:
                            st.download_button("📥 Download Enhanced Presentation", f, file_name="enhanced_presentation.pptx")

            except Exception as e:
                st.error(f"❌ Error: {e}")
